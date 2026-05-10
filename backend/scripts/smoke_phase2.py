from __future__ import annotations

import argparse
import json
import sys
from collections.abc import Iterable
from pathlib import Path
from tempfile import TemporaryDirectory


BACKEND_DIR = Path(__file__).resolve().parents[1]
PROJECT_ROOT = BACKEND_DIR.parent
sys.path.insert(0, str(BACKEND_DIR))

from fastapi.testclient import TestClient  # noqa: E402

from app.core.config import settings  # noqa: E402
from app.main import app  # noqa: E402
from app.services.uploaded_file_parser import parse_uploaded_file  # noqa: E402


def main() -> None:
    parser = argparse.ArgumentParser(description="Phase 02 parser/API smoke tests.")
    parser.add_argument("--keep-output", action="store_true", help="Keep generated data/parsed and data/uploads artifacts.")
    args = parser.parse_args()

    settings.parsed_data_dir.mkdir(parents=True, exist_ok=True)
    settings.upload_dir.mkdir(parents=True, exist_ok=True)
    before_parsed = set(settings.parsed_data_dir.glob("raw_*.json"))
    before_uploads = set(settings.upload_dir.glob("upload_*"))

    try:
        with TemporaryDirectory() as tmp:
            sample_dir = Path(tmp)
            samples = create_samples(sample_dir)
            run_parser_smoke(samples)
            run_api_smoke(samples)
    finally:
        if not args.keep_output:
            cleanup_new_files(settings.parsed_data_dir.glob("raw_*.json"), before_parsed)
            cleanup_new_files(settings.upload_dir.glob("upload_*"), before_uploads)

    print("phase2 smoke ok")


def create_samples(sample_dir: Path) -> dict[str, Path]:
    samples: dict[str, Path] = {}

    txt = sample_dir / "sample.txt"
    txt.write_text("第一章 绪论\n这是第一段内容。\n第二章 方法\n这是第二段内容。", encoding="utf-8")
    samples["txt"] = txt

    md = sample_dir / "sample.md"
    md.write_text("# 第一章 绪论\nMarkdown 内容。\n## 第一节 背景\n更多内容。", encoding="utf-8")
    samples["md"] = md

    csv_path = sample_dir / "sample.csv"
    csv_path.write_text("name,value\nalpha,1\nbeta,2\n", encoding="utf-8")
    samples["csv"] = csv_path

    tsv_path = sample_dir / "sample.tsv"
    tsv_path.write_text("name\tvalue\nalpha\t1\nbeta\t2\n", encoding="utf-8")
    samples["tsv"] = tsv_path

    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "SheetA"
    sheet.append(["topic", "value"])
    sheet.append(["alpha", 1])
    sheet.append(["beta", 2])
    xlsx = sample_dir / "sample.xlsx"
    workbook.save(xlsx)
    samples["xlsx"] = xlsx

    document = Document()
    document.add_heading("第一章 绪论", level=1)
    document.add_paragraph("Word 正文内容。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "值"
    table.cell(1, 0).text = "alpha"
    table.cell(1, 1).text = "1"
    docx = sample_dir / "sample.docx"
    document.save(docx)
    samples["docx"] = docx

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "第一章 绪论"
    slide.placeholders[1].text = "PPT 正文内容。"
    try:
        slide.notes_slide.notes_text_frame.text = "备注内容"
    except Exception:
        pass
    pptx = sample_dir / "sample.pptx"
    presentation.save(pptx)
    samples["pptx"] = pptx

    pdf = PROJECT_ROOT / "第一届AI全栈黑客松赛题.pdf"
    if pdf.exists():
        samples["pdf"] = pdf

    unsupported = sample_dir / "legacy.doc"
    unsupported.write_text("legacy office binary placeholder", encoding="utf-8")
    samples["unsupported"] = unsupported
    return samples


def run_parser_smoke(samples: dict[str, Path]) -> None:
    expected_formats = ["txt", "md", "csv", "tsv", "xlsx", "docx", "pptx"]
    if "pdf" in samples:
        expected_formats.append("pdf")

    for file_format in expected_formats:
        parsed, output_path = parse_uploaded_file(samples[file_format], samples[file_format].name)
        assert output_path.exists(), output_path
        assert parsed.raw_file.format == file_format, (file_format, parsed.raw_file.format)
        assert parsed.elements, file_format
        assert parsed.sections, file_format
        assert parsed.chunks, file_format
        assert_source_locators(file_format, parsed.model_dump(mode="json"))
        print(
            "parser",
            file_format,
            len(parsed.elements),
            len(parsed.sections),
            len(parsed.chunks),
            parsed.metadata["section_strategy"],
        )

    try:
        parse_uploaded_file(samples["unsupported"], samples["unsupported"].name)
    except ValueError as exc:
        assert "LibreOffice" in str(exc)
    else:
        raise AssertionError("unsupported legacy Office file should fail clearly")


def run_api_smoke(samples: dict[str, Path]) -> None:
    client = TestClient(app)

    health = client.get("/api/health")
    assert health.status_code == 200, health.text
    assert health.json()["status"] == "ok"

    with samples["txt"].open("rb") as handle:
        response = client.post("/api/textbooks/upload", files={"file": ("sample.txt", handle, "text/plain")})
    assert response.status_code == 200, response.text
    upload_payload = response.json()
    assert upload_payload["parsed_textbook"]["raw_file"]["format"] == "txt"
    assert_source_locators("txt", upload_payload["parsed_textbook"])

    raw_file_id = upload_payload["raw_file_id"]
    detail = client.get(f"/api/textbooks/{raw_file_id}")
    assert detail.status_code == 200, detail.text

    reparse = client.post(f"/api/textbooks/{raw_file_id}/parse")
    assert reparse.status_code == 200, reparse.text
    assert reparse.json()["raw_file_id"] == raw_file_id

    with samples["unsupported"].open("rb") as bad_single:
        failed_upload = client.post(
            "/api/textbooks/upload",
            files={"file": ("legacy.doc", bad_single, "application/msword")},
        )
    assert failed_upload.status_code == 400, failed_upload.text
    failed_payload = failed_upload.json()
    assert failed_payload["message"] == "教材解析失败", failed_payload
    assert failed_payload["code"] == "PARSE_FAILED", failed_payload
    assert "LibreOffice" in failed_payload["detail"], failed_payload

    with samples["md"].open("rb") as md_handle, samples["unsupported"].open("rb") as bad_handle:
        batch = client.post(
            "/api/textbooks/upload-batch",
            files=[
                ("files", ("sample.md", md_handle, "text/markdown")),
                ("files", ("legacy.doc", bad_handle, "application/msword")),
            ],
        )
    assert batch.status_code == 200, batch.text
    batch_payload = batch.json()
    assert batch_payload["success_count"] == 1, json.dumps(batch_payload, ensure_ascii=False)
    assert batch_payload["failed_count"] == 1, json.dumps(batch_payload, ensure_ascii=False)
    assert batch_payload["items"][0]["parsed_textbook"]["raw_file"]["format"] == "md"
    assert "LibreOffice" in batch_payload["errors"][0]["error"]

    job_id = batch_payload["job"]["id"]
    job = client.get(f"/api/jobs/{job_id}")
    assert job.status_code == 200, job.text
    assert job.json()["status"] == "completed"


def assert_source_locators(file_format: str, parsed_payload: dict[str, object]) -> None:
    chunks = parsed_payload["chunks"]
    assert isinstance(chunks, list) and chunks, file_format
    for chunk in chunks:
        locator = chunk["source_locator"]
        assert locator["raw_file_id"] == parsed_payload["raw_file"]["id"]
        assert locator["source_path"], file_format
        assert locator["locator_text"], file_format
        assert locator["quote_hash"], file_format
        assert locator["char_start"] is not None, file_format
        assert locator["char_end"] is not None, file_format
        assert locator["element_ids"], file_format

    first_locator = chunks[0]["source_locator"]
    if file_format == "pdf":
        assert first_locator["page_start"] is not None, file_format
        assert first_locator["page_end"] is not None, file_format
    elif file_format in {"txt", "md", "docx"}:
        assert first_locator["line_start"] is not None, file_format
        assert first_locator["line_end"] is not None, file_format
    elif file_format in {"csv", "tsv", "xlsx"}:
        assert first_locator["sheet_name"], file_format
        assert first_locator["row_start"] is not None, file_format
        assert first_locator["row_end"] is not None, file_format
    elif file_format == "pptx":
        assert first_locator["slide_number"] is not None, file_format


def cleanup_new_files(paths: Iterable[Path], before: set[Path]) -> None:
    for path in paths:
        if path in before:
            continue
        try:
            path.unlink()
        except FileNotFoundError:
            pass


if __name__ == "__main__":
    main()
