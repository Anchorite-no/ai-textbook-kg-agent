from __future__ import annotations

import csv
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from app.models.schemas import DocumentElementType, ParsedTextbook
from app.services.parser_utils import (
    UnitRecord,
    build_parsed_from_units,
    build_raw_file,
    clean_text,
    save_parsed,
)
from app.services.text_file_parser import SUPPORTED_TEXT_EXTENSIONS, parse_uploaded_text_file


SUPPORTED_UPLOAD_EXTENSIONS = {
    *SUPPORTED_TEXT_EXTENSIONS,
    ".pdf",
    ".docx",
    ".xlsx",
    ".csv",
    ".tsv",
    ".pptx",
}

LEGACY_OFFICE_EXTENSIONS = {
    ".doc": ".docx",
    ".xls": ".xlsx",
    ".ppt": ".pptx",
}

SPREADSHEET_ROW_WINDOW = 40
MAX_TABLE_COLUMNS = 24
PDF_PAGE_WINDOW = 5


@dataclass(frozen=True)
class PdfPageText:
    page: int
    text: str


def parse_uploaded_file(file_path: Path, original_filename: str | None = None) -> tuple[ParsedTextbook, Path]:
    suffix = file_path.suffix.lower()
    if suffix in SUPPORTED_TEXT_EXTENSIONS:
        return parse_uploaded_text_file(file_path, original_filename=original_filename)
    if suffix == ".pdf":
        parsed = _parse_pdf(file_path, original_filename)
    elif suffix == ".docx":
        parsed = _parse_docx(file_path, original_filename)
    elif suffix == ".xlsx":
        parsed = _parse_xlsx(file_path, original_filename)
    elif suffix in {".csv", ".tsv"}:
        parsed = _parse_csv(file_path, original_filename)
    elif suffix == ".pptx":
        parsed = _parse_pptx(file_path, original_filename)
    else:
        if suffix in LEGACY_OFFICE_EXTENSIONS:
            target = LEGACY_OFFICE_EXTENSIONS[suffix]
            raise ValueError(f"旧版 Office 文件 {suffix} 需要先通过 LibreOffice 转换为 {target}，当前阶段未检测到转换器")
        raise ValueError(f"当前计划 02 暂不支持 {suffix or 'unknown'}，请先上传 txt/md/pdf/docx/xlsx/csv/tsv/pptx")
    return parsed, save_parsed(parsed)


def _parse_pdf(file_path: Path, original_filename: str | None) -> ParsedTextbook:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("缺少 pypdf 依赖，请先安装 backend/requirements.txt") from exc

    reader = PdfReader(str(file_path))
    pages: list[PdfPageText] = []
    title = Path(original_filename or file_path.name).stem

    for index, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        pages.append(PdfPageText(page=index, text=text))

    cleaned_pages = _filter_pdf_repeated_headers(pages)
    units = _pdf_sections_from_pages(cleaned_pages, title)
    if not units:
        raise ValueError("PDF 没有可提取文字层，当前阶段暂不做 OCR")

    raw_file = build_raw_file(file_path, original_filename, "pdf", sum(len(page.text) for page in cleaned_pages))
    raw_file = raw_file.model_copy(update={"page_count": len(reader.pages)})
    strategy = "pdf_heading_sections" if any(unit.metadata and unit.metadata.get("detected_heading") for unit in units) else "pdf_page_windows"
    return build_parsed_from_units(raw_file, units, strategy)


def _parse_docx(file_path: Path, original_filename: str | None) -> ParsedTextbook:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("缺少 python-docx 依赖，请先安装 backend/requirements.txt") from exc

    document = Document(str(file_path))
    title = Path(original_filename or file_path.name).stem
    blocks = _docx_blocks(document)
    units: list[UnitRecord] = []
    current_title = "正文"
    current_lines: list[str] = []
    current_start = 1
    unit_index = 1

    def flush(end_line: int) -> None:
        nonlocal unit_index, current_lines, current_start, current_title
        text = clean_text("\n".join(current_lines))
        if not text:
            current_lines = []
            return
        units.append(
            UnitRecord(
                index=unit_index,
                title=current_title,
                text=text,
                element_type=DocumentElementType.paragraph,
                locator_text=f"{title} paragraphs {current_start}-{end_line}",
                line_start=current_start,
                line_end=end_line,
                metadata={"paragraph_start": current_start, "paragraph_end": end_line},
            )
        )
        unit_index += 1
        current_lines = []

    for line_number, kind, text, style in blocks:
        is_heading = kind == "paragraph" and style.lower().startswith("heading")
        if is_heading:
            flush(line_number - 1)
            current_title = clean_text(text) or f"Heading {line_number}"
            current_start = line_number
            continue
        current_lines.append(text)
    flush(blocks[-1][0] if blocks else 1)

    raw_file = build_raw_file(file_path, original_filename, "docx", sum(len(unit.text) for unit in units))
    return build_parsed_from_units(raw_file, units, "docx_headings")


def _docx_blocks(document: object) -> list[tuple[int, str, str, str]]:
    from docx.document import Document as DocxDocument
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    blocks: list[tuple[int, str, str, str]] = []
    line_number = 1
    if not isinstance(document, DocxDocument):
        return blocks
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, document)
            text = clean_text(paragraph.text)
            if text:
                blocks.append((line_number, "paragraph", text, paragraph.style.name if paragraph.style else ""))
                line_number += 1
        elif isinstance(child, CT_Tbl):
            table = Table(child, document)
            table_text = _markdown_table([[cell.text for cell in row.cells] for row in table.rows])
            if table_text:
                blocks.append((line_number, "table", table_text, "Table"))
                line_number += 1
    return blocks


def _parse_xlsx(file_path: Path, original_filename: str | None) -> ParsedTextbook:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("缺少 openpyxl 依赖，请先安装 backend/requirements.txt") from exc

    workbook = load_workbook(str(file_path), read_only=True, data_only=True)
    title = Path(original_filename or file_path.name).stem
    units: list[UnitRecord] = []
    unit_index = 1
    for sheet in workbook.worksheets:
        rows = [(row_index, ["" if value is None else str(value) for value in row]) for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1)]
        non_empty = [(row_index, _trim_empty_tail(values)) for row_index, values in rows if any(str(value).strip() for value in values)]
        if not non_empty:
            continue
        header_row_index, header = _detect_header(non_empty)
        for row_group in _chunk_rows(non_empty, SPREADSHEET_ROW_WINDOW, header_row_index):
            row_start = row_group[0][0]
            row_end = row_group[-1][0]
            table_text = _markdown_table_with_header(row_group, header)
            units.append(
                UnitRecord(
                    index=unit_index,
                    title=f"{sheet.title} rows {row_start}-{row_end}",
                    text=table_text,
                    element_type=DocumentElementType.sheet,
                    locator_text=f"{title} sheet {sheet.title} rows {row_start}-{row_end}",
                    sheet_name=sheet.title,
                    row_start=row_start,
                    row_end=row_end,
                    metadata={
                        "sheet_name": sheet.title,
                        "row_start": row_start,
                        "row_end": row_end,
                        "header_row": header_row_index,
                    },
                )
            )
            unit_index += 1
    workbook.close()

    raw_file = build_raw_file(file_path, original_filename, "xlsx", sum(len(unit.text) for unit in units))
    return build_parsed_from_units(raw_file, units, "xlsx_sheets_rows")


def _parse_csv(file_path: Path, original_filename: str | None) -> ParsedTextbook:
    title = Path(original_filename or file_path.name).stem
    text = _decode_csv(file_path)
    delimiter = "\t" if file_path.suffix.lower() == ".tsv" else _detect_delimiter(text)
    rows = list(csv.reader(text.splitlines(), delimiter=delimiter))
    non_empty = [(index, _trim_empty_tail(row)) for index, row in enumerate(rows, start=1) if any(cell.strip() for cell in row)]
    units: list[UnitRecord] = []
    if not non_empty:
        raise ValueError("CSV/TSV 文件没有可解析的非空行")
    header_row_index, header = _detect_header(non_empty)
    for unit_index, row_group in enumerate(_chunk_rows(non_empty, SPREADSHEET_ROW_WINDOW, header_row_index), start=1):
        row_start = row_group[0][0]
        row_end = row_group[-1][0]
        table_text = _markdown_table_with_header(row_group, header)
        units.append(
            UnitRecord(
                index=unit_index,
                title=f"{file_path.suffix.lower().lstrip('.').upper()} rows {row_start}-{row_end}",
                text=table_text,
                element_type=DocumentElementType.sheet,
                locator_text=f"{title} rows {row_start}-{row_end}",
                sheet_name="CSV",
                row_start=row_start,
                row_end=row_end,
                metadata={
                    "sheet_name": "CSV",
                    "row_start": row_start,
                    "row_end": row_end,
                    "header_row": header_row_index,
                    "delimiter": delimiter,
                },
            )
        )

    file_format = "tsv" if file_path.suffix.lower() == ".tsv" else "csv"
    raw_file = build_raw_file(file_path, original_filename, file_format, sum(len(unit.text) for unit in units))
    return build_parsed_from_units(raw_file, units, f"{file_format}_rows")


def _parse_pptx(file_path: Path, original_filename: str | None) -> ParsedTextbook:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("缺少 python-pptx 依赖，请先安装 backend/requirements.txt") from exc

    presentation = Presentation(str(file_path))
    title = Path(original_filename or file_path.name).stem
    units: list[UnitRecord] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        slide_title = f"Slide {slide_number}"
        for shape in sorted(slide.shapes, key=lambda item: (getattr(item, "top", 0), getattr(item, "left", 0))):
            if getattr(shape, "has_text_frame", False):
                text = clean_text(shape.text)
                if text:
                    if slide_title == f"Slide {slide_number}":
                        slide_title = text.splitlines()[0][:80]
                    slide_lines.append(text)
            if getattr(shape, "has_table", False):
                table_rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                table_text = _markdown_table(table_rows)
                if table_text:
                    slide_lines.append(table_text)
        notes = _extract_slide_notes(slide)
        if notes:
            slide_lines.append(f"Speaker Notes\n{notes}")
        text = clean_text("\n".join(slide_lines))
        if not text:
            continue
        units.append(
            UnitRecord(
                index=slide_number,
                title=slide_title,
                text=text,
                element_type=DocumentElementType.slide,
                locator_text=f"{title} slide {slide_number}",
                slide_number=slide_number,
                metadata={"slide_number": slide_number, "has_notes": bool(notes)},
            )
        )

    raw_file = build_raw_file(file_path, original_filename, "pptx", sum(len(unit.text) for unit in units))
    raw_file = raw_file.model_copy(update={"page_count": len(presentation.slides)})
    return build_parsed_from_units(raw_file, units, "pptx_slides")


def _decode_csv(file_path: Path) -> str:
    raw = file_path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _pdf_sections_from_pages(pages: list[PdfPageText], title: str) -> list[UnitRecord]:
    headings = [_find_pdf_heading(page.text) for page in pages]
    if any(headings):
        return _pdf_heading_sections(pages, headings, title)
    return _pdf_window_sections(pages, title)


def _pdf_heading_sections(pages: list[PdfPageText], headings: list[str | None], title: str) -> list[UnitRecord]:
    units: list[UnitRecord] = []
    current_title: str | None = None
    current_pages: list[PdfPageText] = []

    def flush() -> None:
        nonlocal current_pages, current_title
        if not current_pages:
            return
        page_start = current_pages[0].page
        page_end = current_pages[-1].page
        unit_title = current_title or f"Pages {page_start}-{page_end}"
        units.append(
            UnitRecord(
                index=len(units) + 1,
                title=unit_title,
                text=_join_pdf_pages(current_pages),
                element_type=DocumentElementType.page,
                locator_text=f"{title} pages {page_start}-{page_end}",
                page=page_start,
                page_end=page_end,
                metadata={
                    "page_start": page_start,
                    "page_end": page_end,
                    "detected_heading": current_title,
                },
            )
        )
        current_pages = []

    for page, heading in zip(pages, headings, strict=False):
        if heading and current_pages:
            flush()
            current_title = heading
        elif heading and not current_pages:
            current_title = heading
        current_pages.append(page)
    flush()
    return units


def _pdf_window_sections(pages: list[PdfPageText], title: str) -> list[UnitRecord]:
    units: list[UnitRecord] = []
    non_empty = [page for page in pages if page.text.strip()]
    for index in range(0, len(non_empty), PDF_PAGE_WINDOW):
        group = non_empty[index : index + PDF_PAGE_WINDOW]
        page_start = group[0].page
        page_end = group[-1].page
        units.append(
            UnitRecord(
                index=len(units) + 1,
                title=f"Pages {page_start}-{page_end}",
                text=_join_pdf_pages(group),
                element_type=DocumentElementType.page,
                locator_text=f"{title} pages {page_start}-{page_end}",
                page=page_start,
                page_end=page_end,
                metadata={"page_start": page_start, "page_end": page_end},
            )
        )
    return units


def _join_pdf_pages(pages: list[PdfPageText]) -> str:
    return "\n\n".join(f"第 {page.page} 页\n{page.text}" for page in pages if page.text.strip())


def _find_pdf_heading(text: str) -> str | None:
    for line in text.splitlines()[:10]:
        candidate = clean_text(line)
        if len(candidate) <= 48 and re.match(r"^(绪论|第[一二三四五六七八九十百0-9]+[章节篇])", candidate):
            return candidate
    return None


def _filter_pdf_repeated_headers(pages: list[PdfPageText]) -> list[PdfPageText]:
    edge_counter: dict[str, int] = {}
    for page in pages:
        lines = [clean_text(line) for line in page.text.splitlines() if clean_text(line)]
        for line in [*lines[:2], *lines[-2:]]:
            if len(line) <= 40:
                edge_counter[line] = edge_counter.get(line, 0) + 1
    threshold = max(3, len(pages) // 4)
    repeated = {line for line, count in edge_counter.items() if count >= threshold}

    cleaned: list[PdfPageText] = []
    for page in pages:
        lines = [clean_text(line) for line in page.text.splitlines()]
        kept = [
            line
            for index, line in enumerate(lines)
            if line and not (line in repeated and (index <= 1 or index >= len(lines) - 2))
        ]
        cleaned.append(PdfPageText(page=page.page, text="\n".join(kept)))
    return cleaned


def _chunk_rows(
    rows: list[tuple[int, list[str]]],
    size: int,
    header_row_index: int | None = None,
) -> Iterable[list[tuple[int, list[str]]]]:
    data_rows = [row for row in rows if row[0] != header_row_index]
    if not data_rows and rows:
        yield rows
        return
    for index in range(0, len(data_rows), size):
        yield data_rows[index : index + size]


def _detect_header(rows: list[tuple[int, list[str]]]) -> tuple[int | None, list[str] | None]:
    if not rows:
        return None, None
    row_index, values = rows[0]
    non_empty = [value for value in values if value.strip()]
    if len(non_empty) >= 2:
        return row_index, values
    return None, None


def _trim_empty_tail(values: list[str]) -> list[str]:
    trimmed = list(values)
    while trimmed and not str(trimmed[-1]).strip():
        trimmed.pop()
    trimmed = [clean_text(str(value)) for value in trimmed]
    if len(trimmed) > MAX_TABLE_COLUMNS:
        return [*trimmed[:MAX_TABLE_COLUMNS], f"...({len(trimmed) - MAX_TABLE_COLUMNS} more columns)"]
    return trimmed


def _markdown_table(rows: list[list[str]]) -> str:
    clean_rows = [[clean_text(cell) for cell in row] for row in rows if any(clean_text(cell) for cell in row)]
    if not clean_rows:
        return ""
    width = max(len(row) for row in clean_rows)
    padded = [row + [""] * (width - len(row)) for row in clean_rows]
    header = padded[0]
    body = padded[1:] or [[""] * width]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _markdown_table_with_header(
    row_group: list[tuple[int, list[str]]],
    header: list[str] | None,
) -> str:
    rows = [values for _row_index, values in row_group]
    if header:
        rows = [header, *rows]
    return _markdown_table(rows)


def _detect_delimiter(text: str) -> str:
    try:
        dialect = csv.Sniffer().sniff(text[:4096], delimiters=",;\t|")
        return dialect.delimiter
    except csv.Error:
        return ","


def _extract_slide_notes(slide: object) -> str:
    try:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        return clean_text(text_frame.text or "")
    except Exception:  # noqa: BLE001 - notes are optional and backend parsing should continue.
        return ""
